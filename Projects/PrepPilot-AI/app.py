import streamlit as st
import os
from PyPDF2 import PdfReader
import math
import re
from datetime import datetime, date
import io
from utils.ai import summarize_text, search_subject_info, generate_study_map, detect_modules
from pptx import Presentation
from utils.pdf_export import generate_pdf

with open("assets/style.css") as f:
    st.markdown(
        f"<style>{f.read()}</style>",
        unsafe_allow_html=True
    )
st.set_page_config(
    page_title="PrepPilot AI",
    page_icon="📚",
    layout="wide"
)

st.title("📚 PrepPilot AI")
st.subheader("Your AI-Powered Study Assistant")

st.markdown("---")

st.write("""
Welcome to **PrepPilot AI**!

Upload your study materials and let AI help you build a winning exam strategy.
""")

st.markdown("### Step 1: Upload your study materials")
uploaded_files = st.file_uploader(
    "📄 Upload one or more PDFs or PowerPoint files (any modules you have notes for)",
    type=["pdf", "pptx"],
    accept_multiple_files=True
)

st.markdown("### Step 2: Subject details")
subject_code = st.text_input(
    "📘 KTU Subject Code",
    placeholder="Example: CST301"
).strip().upper()

exam_date = st.date_input("📅 Exam Date")

generate_clicked = st.button("🎯 Generate Study Map")
@st.cache_data
def extract_pdf_text(file_bytes, file_name):
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text, len(reader.pages)

@st.cache_data
def extract_pptx_text(file_bytes, file_name):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
                text += "\n"
    return text, slide_count


if generate_clicked:
    if not uploaded_files:
        st.error("Please upload at least one PDF before generating a study map.")
    elif not subject_code:
        st.error("Please enter your KTU subject code.")
    else:
        combined_text = ""
        file_snippets = {}

        for uploaded_file in uploaded_files:
            file_bytes = uploaded_file.getvalue()

            if uploaded_file.name.lower().endswith(".pptx"):
                text, unit_count = extract_pptx_text(file_bytes, uploaded_file.name)
            else:
                text, unit_count = extract_pdf_text(file_bytes, uploaded_file.name)

            combined_text += f"\n\n--- Content from {uploaded_file.name} ---\n{text}"
            file_snippets[uploaded_file.name] = text[:800]

        days_left = (exam_date - date.today()).days
        if days_left < 0:
            st.warning("The exam date you entered is in the past. Please check the date.")

        with st.spinner(f"Step 1/3: Researching {subject_code} syllabus..."):
            subject_info = search_subject_info(subject_code)

        with st.spinner("Step 2/3: Matching your files to modules..."):
            module_mapping = detect_modules(file_snippets, subject_info)

        st.markdown("### 📊 Your Files, Mapped to Modules")
        st.markdown(module_mapping)

        with st.spinner("Step 3/3: Building your personalized study map..."):
            study_data = generate_study_map(
                combined_text=combined_text,
                subject_code=subject_code,
                subject_info=subject_info,
                days_left=days_left
            )

        with st.container(border=True):

         st.markdown("---")
         st.info(f"⏳ {days_left} day(s) left until your exam.")
         st.markdown("## 🗺️ Your Study Map")

         # --- Pie chart: study order / difficulty breakdown ---
         modules = study_data["modules"]
         names = [m["name"] for m in modules]
         scores = [m["difficulty_score"] for m in modules]

         col_chart, col_order = st.columns([1, 1])

         with col_chart:
            st.markdown("#### Difficulty Breakdown")
            st.plotly_chart(
                {
                    "data": [{
                        "labels": names,
                        "values": scores,
                        "type": "pie",
                        "hole": 0.4
                    }],
                    "layout": {"showlegend": True}
                },
                use_container_width=True
            )

         with col_order:
            st.markdown("#### Recommended Study Order")
            sorted_modules = sorted(modules, key=lambda m: m["study_order"])
            for i, m in enumerate(sorted_modules, 1):
                st.markdown(f"**{i}.** {m['name']} (Difficulty: {m['difficulty_score']}/10)")

         st.markdown("---")

         # --- Tabs: one per module ---
         tab_labels = [m["name"] for m in modules]
         tabs = st.tabs(tab_labels)

         for tab, module in zip(tabs, modules):
            with tab:
                st.markdown("**Key Points:**")
                for point in module["key_points"]:
                    st.markdown(f"- {point}")

                st.markdown("**Likely Questions:**")
                for q in module["likely_questions"]:
                    with st.container():
                        st.markdown(f"**Q: {q['question']}** _({q['marks']} marks)_")

                        # Split model_answer into text and code parts
                        answer_parts = q["model_answer"].split("```")
                        for i, part in enumerate(answer_parts):
                            if i % 2 == 1:
                                # Odd-indexed parts are inside ``` fences = code
                                lines = part.strip().split("\n")
                                lang = lines[0] if lines[0].isalpha() else ""
                                code_content = "\n".join(lines[1:]) if lang else part.strip()
                                st.code(code_content, language=lang if lang else None)
                            else:
                                if part.strip():
                                    st.markdown(part.strip())

                        if q.get("needs_diagram"):
                            search_url = f"https://www.google.com/search?q={q['diagram_search_term'].replace(' ', '+')}&tbm=isch"
                            st.markdown(f"[🖼️ View diagram reference for this answer]({search_url})")
                        st.markdown("---")
        st.markdown("---")
        pdf_bytes = generate_pdf(study_data, subject_code, days_left)
        st.download_button(
         label="📥 Download Study Map as PDF",
         data=pdf_bytes,
         file_name=f"PrepPilot_{subject_code}_StudyMap.pdf",
         mime="application/pdf"
         )

st.sidebar.title("PrepPilot AI")
st.sidebar.success("Version 2.0")