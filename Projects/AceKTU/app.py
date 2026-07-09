import streamlit as st
import os
from PyPDF2 import PdfReader
import math
import re
from datetime import datetime, date
import io
from utils.ai import summarize_text, search_subject_info, generate_study_map, detect_modules
from pptx import Presentation
from utils.pdf_export import generate_pdf
import pathlib
st.markdown(
    """
    <link href="https://fonts.googleapis.com/css2?family=Orbitron:wght@500;700;900&family=Rajdhani:wght@400;500;600&display=swap" rel="stylesheet">
    """,
    unsafe_allow_html=True
)

css_path = pathlib.Path(__file__).parent / "assets" / "style.css"
with open(css_path) as f:
    st.markdown(
        f"<style>{f.read()}</style>",
        unsafe_allow_html=True
    )
st.set_page_config(
    page_title="AceKTU",
    page_icon="📚",
    layout="wide"
)

st.title("📚 AceKTU")
st.subheader("Your AI-Powered Study Assistant")

st.markdown("---")

st.markdown(
    """
    <div style="text-align:center;">
        <p style="font-size:20px;">Welcome to <strong>AceKTU</strong>!</p>
        <p style="font-size:16px; color:#8fa3c0;">Let's make your exam study map</p>
    </div>
    """,
    unsafe_allow_html=True
)
st.markdown(
    """
    <div style="display:flex; justify-content:center; gap:40px; margin: 30px 0; text-align:center;">
        <div><div style="font-size:28px;">📤</div><div style="color:#8fa3c0; font-size:13px;">1. Upload Notes</div></div>
        <div><div style="font-size:28px;">🔍</div><div style="color:#8fa3c0; font-size:13px;">2. AI Researches</div></div>
        <div><div style="font-size:28px;">🗺️</div><div style="color:#8fa3c0; font-size:13px;">3. Get Your Plan</div></div>
    </div>
    """,
    unsafe_allow_html=True
)
col_left, col_center, col_right = st.columns([1, 3, 1])
with col_center:
    st.markdown(
        """
        <div style="text-align:center; margin-bottom: 10px;">
            <div style="font-size: 48px;">📤</div>
            <div style="font-size: 22px; font-weight: 700; color: #94a3b8;">Upload Your Study Materials</div>
            <div style="font-size: 14px; color: #8fa3c0;">PDF or PowerPoint — any modules you have notes for</div>
        </div>
        """,
        unsafe_allow_html=True
    )
    uploaded_files = st.file_uploader(
        "Upload files",
        type=["pdf", "pptx"],
        accept_multiple_files=True,
        label_visibility="collapsed"
    )


col_subject, col_date = st.columns(2)

with col_subject:
    subject_code = st.text_input(
        "📘 KTU Subject Code",
        placeholder="Example: CST301"
    ).strip().upper()

with col_date:
    exam_date = st.date_input("📅 Exam Date")

generate_clicked = st.button("🎯 Generate Study Map", key="generate_main_btn")
if not generate_clicked and "study_data" not in st.session_state:
    st.markdown(
        """
        <div style="text-align:center; padding: 40px 0; color:#5a6d8a;">
            <div style="font-size:40px;">🗺️</div>
            <div style="font-size:16px;">Your personalized study map will appear here</div>
        </div>
        """,
        unsafe_allow_html=True
    )


@st.cache_data
def extract_pdf_text(file_bytes, file_name):
    reader = PdfReader(io.BytesIO(file_bytes))
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text, len(reader.pages)


@st.cache_data
def extract_pptx_text(file_bytes, file_name):
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + " "
                text += "\n"
    return text, slide_count


if generate_clicked:
    if not uploaded_files:
        st.error("Please upload at least one PDF before generating a study map.")
    elif not subject_code:
        st.error("Please enter your KTU subject code.")
    else:
        combined_text = ""
        file_snippets = {}

        for uploaded_file in uploaded_files:
            file_bytes = uploaded_file.getvalue()

            if uploaded_file.name.lower().endswith(".pptx"):
                text, unit_count = extract_pptx_text(file_bytes, uploaded_file.name)
            else:
                text, unit_count = extract_pdf_text(file_bytes, uploaded_file.name)

            combined_text += f"\n\n--- Content from {uploaded_file.name} ---\n{text}"
            file_snippets[uploaded_file.name] = text[:800]

        days_left = (exam_date - date.today()).days
        if days_left < 0:
            st.warning("The exam date you entered is in the past. Please check the date.")
            st.stop()

        try:
            with st.spinner(f"🔍 Step 1/3: Researching {subject_code} syllabus..."):
                subject_info = search_subject_info(subject_code)

            if "SUBJECT_NOT_FOUND" in subject_info:
                st.error(f"⚠️ '{subject_code}' doesn't appear to be a valid KTU subject code. Please check and try again.")
                st.stop()

            with st.spinner("📋 Step 2/3: Matching your files to modules..."):
                module_mapping = detect_modules(file_snippets, subject_info)

            st.markdown("### 📊 Your Files, Mapped to Modules")
            st.markdown(module_mapping)
        except Exception as e:
            error_message = str(e)
            if "RESOURCE_EXHAUSTED" in error_message or "429" in error_message:
                st.error("⚠️ Daily AI request limit reached (free tier allows 20/day). Please try again tomorrow, or consider upgrading your Gemini plan.")
            elif "UNAVAILABLE" in error_message or "503" in error_message:
                st.error("⚠️ Gemini's servers are temporarily busy. Please wait a moment and try again.")
            else:
                print(f"--- UNEXPECTED ERROR ---\n{type(e).__name__}: {e}")
                st.error(f"⚠️ Something went wrong: {type(e).__name__}. Check your terminal for details.")
            st.stop()

        try:
            with st.spinner("🗺️ Step 3/3: Building your personalized study map..."):
                study_data = generate_study_map(
                    combined_text=combined_text,
                    subject_code=subject_code,
                    subject_info=subject_info,
                    days_left=days_left
                )
        except Exception as e:
            error_message = str(e)
            if "RESOURCE_EXHAUSTED" in error_message or "429" in error_message:
                st.error("⚠️ Daily AI request limit reached (free tier allows 20/day). Please try again tomorrow, or consider upgrading your Gemini plan.")
            elif "UNAVAILABLE" in error_message or "503" in error_message:
                st.error("⚠️ Gemini's servers are temporarily busy. Please wait a moment and try again.")
            else:
                print(f"--- UNEXPECTED ERROR ---\n{type(e).__name__}: {e}")
                st.error(f"⚠️ Something went wrong: {type(e).__name__}. Check your terminal for details.")
            st.stop()

        st.session_state["study_data"] = study_data
        st.session_state["days_left"] = days_left
        st.session_state["subject_code"] = subject_code


if "study_data" in st.session_state:
    study_data = st.session_state["study_data"]
    days_left = st.session_state["days_left"]
    subject_code = st.session_state["subject_code"]
    modules = study_data["modules"]

    with st.container(border=True, key="study_map_box"):

        st.markdown("---")
        st.info(f"⏳ {days_left} day(s) left until your exam.")
        st.markdown("## 🗺️ Your Study Map")

        if "study_schedule" in study_data:
            st.markdown("### 📅 Your Day-by-Day Plan")
            for entry in study_data["study_schedule"]:
                st.markdown(f"**{entry['days']}:** {entry['focus']}  \n_{entry['reason']}_")
            st.markdown("---")

        names = [m["name"] for m in modules]
        scores = [m["difficulty_score"] for m in modules]

        col_chart, col_order = st.columns([1, 1])

        with col_chart:
            st.markdown("#### Difficulty Breakdown")
            st.plotly_chart(
                {
                    "data": [{
                        "labels": names,
                        "values": scores,
                        "type": "pie",
                        "hole": 0.4
                    }],
                    "layout": {
                        "showlegend": True,
                        "paper_bgcolor": "rgba(0,0,0,0)",
                        "plot_bgcolor": "rgba(0,0,0,0)",
                        "font": {"color": "#ffffff"}
                    }
                },
                use_container_width=True
            )

        with col_order:
            st.markdown("#### Recommended Study Order")
            sorted_modules = sorted(modules, key=lambda m: m["study_order"])
            for i, m in enumerate(sorted_modules, 1):
                st.markdown(f"**{i}.** {m['name']} (Difficulty: {m['difficulty_score']}/10)")
                if "difficulty_reason" in m:
                    st.caption(m["difficulty_reason"])

        st.markdown("---")
        st.markdown("#### 📚 Module Details")

        for idx, module in enumerate(modules):
            module_key = f"module_{idx}"

            if f"{module_key}_open" not in st.session_state:
                st.session_state[f"{module_key}_open"] = False

            if "completed_points" not in st.session_state:
                st.session_state["completed_points"] = {}

            total_points = len(module["key_points"])
            checked_count = sum(
                1 for p_idx in range(total_points)
                if st.session_state["completed_points"].get(f"{module_key}_point_{p_idx}", False)
            )
            completion_pct = int((checked_count / total_points) * 100) if total_points > 0 else 0

            with st.container(key=f"{module_key}_row"):
                col_card, col_ring = st.columns([6, 1])

                with col_card:
                    icon = "▼" if st.session_state[f"{module_key}_open"] else "▶"
                    if st.button(f"{icon} {module['name']}", key=f"{module_key}_toggle", use_container_width=True):
                        st.session_state[f"{module_key}_open"] = not st.session_state[f"{module_key}_open"]

                with col_ring:
                    st.plotly_chart(
                        {
                            "data": [{
                                "values": [completion_pct, 100 - completion_pct],
                                "labels": ["Done", "Remaining"],
                                "type": "pie",
                                "hole": 0.65,
                                "marker": {"colors": ["#94a3b8", "#3a3a3a"]},
                                "textinfo": "none",
                                "showlegend": False,
                                "sort": False,
                                "direction": "clockwise"
                            }],
                            "layout": {
                                "margin": {"l": 10, "r": 10, "t": 10, "b": 10},
                                "height": 90,
                                "paper_bgcolor": "rgba(0,0,0,0)",
                                "plot_bgcolor": "rgba(0,0,0,0)",
                                "annotations": [{
                                    "text": f"{completion_pct}%",
                                    "showarrow": False,
                                    "font": {"size": 13, "color": "#ffffff"}
                                }]
                            }
                        },
                        use_container_width=True,
                        key=f"{module_key}_ring"
                    )

                if st.session_state[f"{module_key}_open"]:
                    with st.container(border=True, key=f"{module_key}_details"):
                        st.markdown("**Key Points:**")
                        def update_completion(dict_key):
                            st.session_state["completed_points"][dict_key] = st.session_state[dict_key]

                        for p_idx, point_data in enumerate(module["key_points"]):
                            point_key = f"{module_key}_point_{p_idx}"
                            st.checkbox(
                                point_data["point"],
                                value=st.session_state["completed_points"].get(point_key, False),
                                key=point_key,
                                on_change=update_completion,
                                args=(point_key,)
                            )
                            st.caption(point_data["content"])
                            study_link_query = f"{point_data['point']} {subject_code} KTU notes"
                            study_link = f"https://www.google.com/search?q={study_link_query.replace(' ', '+')}"
                            st.markdown(f"[🔗 Study this topic in detail]({study_link})")
                            st.markdown("")

                        st.markdown("**Likely Questions:**")
                        for q in module["likely_questions"]:
                            with st.container():
                                st.markdown(f"**Q: {q['question']}** _({q['marks']} marks)_")

                                answer_parts = q["model_answer"].split("```")
                                for i, part in enumerate(answer_parts):
                                    if i % 2 == 1:
                                        lines = part.strip().split("\n")
                                        lang = lines[0] if lines[0].isalpha() else ""
                                        code_content = "\n".join(lines[1:]) if lang else part.strip()
                                        st.code(code_content, language=lang if lang else None)
                                    else:
                                        if part.strip():
                                            st.markdown(part.strip())

                                if q.get("needs_diagram"):
                                    search_url = f"https://www.google.com/search?q={q['diagram_search_term'].replace(' ', '+')}&tbm=isch"
                                    st.markdown(f"[🖼️ View diagram reference for this answer]({search_url})")
                                st.markdown("---")

            st.markdown("")

        st.markdown("---")
        pdf_bytes = generate_pdf(study_data, subject_code, days_left)
        st.download_button(
            label="📥 Download Quick Revision Sheet (PDF)",
            data=pdf_bytes,
            file_name=f"AceKTU_{subject_code}_QuickRevision.pdf",
            mime="application/pdf"
        )


st.sidebar.title("AceKTU")
st.sidebar.success("Version 2.0")
st.markdown("---")
st.markdown(
    """
    <div style="text-align:center; padding: 20px 0; color:#5a6d8a; font-size:13px;">
        Built by <strong style="color:#94a3b8;">Fazil Firose Ibrahim</strong><br>
        <a href="https://github.com/FazilFirose" style="color:#8fa3c0;" target="_blank">GitHub</a> ·
        <a href="https://www.linkedin.com/in/fazil-firose-ibrahim-97a378255" style="color:#8fa3c0;" target="_blank">LinkedIn</a>
    </div>
    """,
    unsafe_allow_html=True
)